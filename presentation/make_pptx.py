#!/usr/bin/env python3
"""Build jaseci_deck.pptx (16:9, native editable text) from the deck content.
Text is real text boxes (editable in Google Slides). TikZ diagrams + result
charts go in as images (figs/). Run: python3 make_pptx.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figs")

INK   = RGBColor(0x1A, 0x1A, 0x1A)
SOFT  = RGBColor(0x73, 0x73, 0x73)
FILL  = RGBColor(0xF4, 0xF4, 0xF4)
BORD  = RGBColor(0xC8, 0xC8, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide():
    return prs.slides.add_slide(BLANK)


def _runs(p, text, size, color, bold_default=False, italic_default=False, mono=False):
    # parse **bold** and *italic* into runs
    import re
    # split keeping markers
    tokens = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*)', text)
    for tok in tokens:
        if tok == "":
            continue
        b, i = bold_default, italic_default
        s = tok
        if tok.startswith("**") and tok.endswith("**"):
            s = tok[2:-2]; b = True
        elif tok.startswith("*") and tok.endswith("*"):
            s = tok[1:-1]; i = True
        r = p.add_run(); r.text = s
        f = r.font
        f.size = Pt(size); f.bold = b; f.italic = i; f.color.rgb = color
        f.name = "Consolas" if mono else "Calibri"


def tbox(s, left, top, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def title(s, text):
    _, tf = tbox(s, 0.5, 0.28, 12.3, 0.8)
    p = tf.paragraphs[0]
    _runs(p, text, 26, INK, bold_default=True)


def para(tf, text, size=15, color=INK, align=PP_ALIGN.LEFT, space=6, first=False,
         bold=False, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space); p.line_spacing = 1.04
    _runs(p, text, size, color, bold_default=bold, italic_default=italic)
    return p


def bullets(s, items, left, top, w, h, size=15, marker="•  ", color=INK, space=7):
    _, tf = tbox(s, left, top, w, h)
    for idx, it in enumerate(items):
        para(tf, marker + it, size=size, color=color, space=space, first=(idx == 0))
    return tf


def numbered(s, items, left, top, w, h, size=15, space=8):
    _, tf = tbox(s, left, top, w, h)
    for idx, it in enumerate(items):
        para(tf, f"{idx+1}.  " + it, size=size, space=space, first=(idx == 0))
    return tf


def code(s, text, left, top, w, h, size=10.5):
    tb, tf = tbox(s, left, top, w, h)
    tf.word_wrap = False
    tb.fill.solid(); tb.fill.fore_color.rgb = FILL
    tb.line.color.rgb = BORD; tb.line.width = Pt(0.75)
    for idx, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(0); p.line_spacing = 1.0
        r = p.add_run(); r.text = ln if ln else " "
        f = r.font; f.size = Pt(size); f.name = "Consolas"; f.color.rgb = INK
    return tb


def caption(s, text, top, size=11, left=0.6, w=12.1, color=SOFT, align=PP_ALIGN.CENTER):
    _, tf = tbox(s, left, top, w, 0.4)
    para(tf, text, size=size, color=color, align=align, first=True, space=0)


def image(s, name, left, top, w=None, h=None, center_x=None):
    path = os.path.join(FIG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w is not None and h is None:
        h = w / ar
    elif h is not None and w is None:
        w = h * ar
    if center_x is not None:
        left = center_x - w / 2
    s.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return w, h


# ----------------------------------------------------------------- 1 TITLE
s = slide()
_, tf = tbox(s, 1.0, 2.1, 11.3, 2.2, anchor=MSO_ANCHOR.TOP)
para(tf, "Teaching a model to write Jac", size=40, color=INK, bold=True, first=True, space=8, align=PP_ALIGN.CENTER)
para(tf, "Synthetic Python → Jac data and a small finetune probe", size=20, color=SOFT, align=PP_ALIGN.CENTER, space=18)
para(tf, "Ayush Madhav Kumar", size=18, color=INK, align=PP_ALIGN.CENTER, space=4)
para(tf, "University of Michigan  •  Jaseci", size=14, color=SOFT, align=PP_ALIGN.CENTER, space=0)
caption(s, "Synthetic Python → Jac   •   MLX LoRA on Apple Silicon   •   Qwen3 Coder & Gemma 4", 6.4, size=12)

# ----------------------------------------------------------------- 2 NO CORPUS
s = slide(); title(s, "No real Jac corpus exists")
bullets(s, [
    "Models trained on Python have weak priors on Jac's **data spatial idioms**: walkers, nodes and edges instead of functions and classes.",
    "Stock 30B and 26B Gemma 4 and Qwen3 Coder produce **0% runnable Jac**. They have never seen the language.",
], 0.6, 1.2, 12.1, 1.5, size=15)
code(s, "# a graph walk written as a loop\n"
        "def traverse(graph, start):\n"
        "    stack = [start]\n"
        "    seen = set()\n"
        "    while stack:\n"
        "        n = stack.pop()\n"
        "        seen.add(n)\n"
        "        stack += graph[n]\n"
        "    return seen", 0.6, 2.9, 6.0, 3.6, size=11)
code(s, "# the graph walks itself\n"
        "walker Traverse {\n"
        "    can step with Item entry {\n"
        "        here.seen = True;\n"
        "        visit [-->];\n"
        "    }\n"
        "}\n"
        "node Item { has seen: bool = False; }", 6.9, 2.9, 5.8, 3.6, size=11)

# ----------------------------------------------------------------- 3 ANCHORS
s = slide(); title(s, "Three anchors for 100% synthetic data")
numbered(s, [
    "**Grammar** is the distribution anchor. The Jac grammar defines the space of valid programs, and every sample is shaped to fit it.",
    "**Compiler and tests** are a free oracle. Labelling is unlimited and costs nothing: reject anything that does not compile, run, or match its tests.",
    "**Python** is the proxy distribution. Translate validated Python into Jac to borrow a large, correct source.",
], 0.6, 1.3, 7.0, 3.8, size=15)
_, tf = tbox(s, 0.6, 5.2, 7.0, 1.4)
para(tf, "The quality gate is jac run: compile it, run it, and the output has to match. Not the type checker, which rejects plenty of code that runs.",
     size=13, color=INK, first=True)
image(s, "diagram-1.png", None, 1.7, w=4.4, center_x=10.3)
caption(s, "the multi PLT method", 5.15, size=12, left=8.0, w=4.6)

# ----------------------------------------------------------------- 4 DATA TIERS PT1
s = slide(); title(s, "Data tiers, part 1: the three sources")
numbered(s, [
    "**Idiomatic core** [147]. Jac written by hand or with an agent (Claude Code plus the Jac MCP): walkers, nodes, edges. *31 of these are graph shaped.*",
    "**Transpile volume** [1500]. Python auto transpiled with jac py2jac and gated on behavior. Correct, but *Python shaped*. This is the cheap volume.",
    "**DPO pairs** [147]. chosen is idiomatic, rejected is the transpile. This teaches the model *what is right versus what is wrong*.",
], 0.6, 2.0, 7.4, 3.6, size=15)
image(s, "diagram-2.png", None, 1.2, w=4.2, center_x=10.4)
caption(s, "a DPO pair is one task with two answers (chosen ≻ rejected)", 3.3, size=11, left=8.2, w=4.4)
caption(s, "**1647** SFT examples (147 idiom, 1500 transpile)   •   **147** DPO pairs   •   gate is jac run", 6.7, size=13, left=0.6, w=12.1, color=INK)

# ----------------------------------------------------------------- 5 DATA TIERS PT2
s = slide(); title(s, "Data tiers, part 2: how the tiers are built")
image(s, "diagram-3.png", None, 2.4, w=12.4, center_x=SW/2)
caption(s, "jac run gated   •   manifest 1:3 idiom:transpile   •   holdout decontaminated (disjoint offsets, 14 token shingles)",
        5.7, size=13, color=INK)

# ----------------------------------------------------------------- 6 EXAMPLES 1/2
s = slide(); title(s, "The data, tier by tier: real examples (1 of 2)")
_, tf = tbox(s, 0.6, 1.15, 6.0, 0.7)
para(tf, "Tier 1  •  Idiomatic core", size=14, bold=True, first=True, space=0)
para(tf, "task: product_of_tree, multiply values over a tree (walker / node / edge)", size=11, color=SOFT, space=0)
code(s, "node TNode { has val: int; }\n\n"
        "walker Product {\n"
        "    has result: int = 1;\n"
        "    can mul with TNode entry {\n"
        "        self.result *= here.val;\n"
        "        visit [-->];   # the walk IS the program\n"
        "    }\n"
        "}\n\n"
        "with entry {\n"
        "    a = TNode(val=2); b = TNode(val=3);\n"
        "    root ++> a;  a ++> b;  a ++> c;\n"
        "    print((a spawn Product()).result);\n"
        "}", 0.6, 1.95, 6.0, 4.0, size=9.5)
_, tf = tbox(s, 6.9, 1.15, 5.8, 0.7)
para(tf, "Tier 2  •  Transpile volume", size=14, bold=True, first=True, space=0)
para(tf, "task: is_even, mechanical, Python shaped (jac py2jac, gated)", size=11, color=SOFT, space=0)
code(s, "def is_even(num: Any) -> object {\n"
        "    return ((num % 2) == 0);\n"
        "}\n\n"
        "with entry {\n"
        "    print(is_even(5));\n"
        "    print(is_even(10));\n"
        "}", 6.9, 1.95, 5.8, 2.4, size=9.5)
_, tf = tbox(s, 6.9, 4.6, 5.8, 1.6)
para(tf, "Correct and runnable, but it is just Python with braces: Any and object types, no walker. The transpile gives cheap scale. The idiomatic core gives the shape.",
     size=12, color=INK, first=True)

# ----------------------------------------------------------------- 7 DPO PAIR 2/2
s = slide(); title(s, "The data, tier by tier: a real DPO pair (2 of 2)")
caption(s, "task: count_children(graph, node), count a node's direct children", 1.2, size=12, color=INK)
_, tf = tbox(s, 0.6, 1.7, 6.0, 0.4)
para(tf, "chosen  (idiomatic, reward up)", size=13, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
code(s, "node GNode { has name: str; }\n\n"
        "walker CountChildren {\n"
        "    has count: int = 0;\n"
        "    can chk with GNode entry {\n"
        "        self.count = len([-->]);\n"
        "        disengage;\n"
        "    }\n"
        "}", 0.6, 2.15, 6.0, 2.8, size=10.5)
_, tf = tbox(s, 6.9, 1.7, 5.8, 0.4)
para(tf, "rejected  (transpile, Python shaped, reward down)", size=13, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
code(s, '"""Return the number of direct\n'
        '   children of node."""\n'
        "def count_children(graph: dict,\n"
        "                   `node: str) -> int {\n"
        "    return len(graph.get(`node, []));\n"
        "}", 6.9, 2.15, 5.8, 2.0, size=10.5)
_, tf = tbox(s, 0.6, 5.4, 12.1, 1.2)
para(tf, "Same task, two correct answers. DPO trains chosen ≻ rejected, pushing generation off graph.get(...) and toward len([-->]).",
     size=14, color=INK, first=True, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------- 8 SFT vs DPO
s = slide(); title(s, "What is SFT and DPO")
bullets(s, [
    "**SFT** is supervised finetuning. Imitate gold outputs, which teaches the model to produce runnable Jac.",
    "**DPO** is direct preference optimization. Learn from chosen ≻ rejected pairs, which sharpens the model toward idiom without a reward model.",
], 0.6, 1.5, 5.6, 3.0, size=16)
_, tf = tbox(s, 0.6, 4.8, 5.6, 1.0)
para(tf, "Here, SFT teaches it to run, and DPO pushes it from runs to idiomatic.", size=14, color=INK, first=True)
image(s, "diagram-4.png", None, 3.0, w=6.4, center_x=9.7)

# ----------------------------------------------------------------- 9 FUNCTION vs GRAPH
s = slide(); title(s, "Function vs graph tasks")
_, tf = tbox(s, 0.6, 1.2, 6.0, 0.7)
para(tf, "Function task", size=16, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
para(tf, "standalone pure function → Jac def", size=11, color=SOFT, space=0, align=PP_ALIGN.CENTER)
image(s, "diagram-5.png", None, 1.95, w=3.6, center_x=3.5)
bullets(s, [
    "idiomatic ≈ transpile → **sim 0.97**",
    "**no idiom headroom**: a pure function legitimately needs no walker",
    "oracle: recorded test cases (150 holdout)",
], 0.6, 3.1, 6.0, 2.4, size=13)
_, tf = tbox(s, 6.9, 1.2, 5.8, 0.7)
para(tf, "Graph task", size=16, bold=True, first=True, space=0, align=PP_ALIGN.CENTER)
para(tf, "dict and stack traversal → nodes / edges / walker", size=11, color=SOFT, space=0, align=PP_ALIGN.CENTER)
image(s, "diagram-6.png", None, 1.95, w=3.8, center_x=9.8)
bullets(s, [
    "idiomatic *diverges* hard → **sim 0.26**, ~8 constructs",
    "**real idiom headroom**: this is where DPO can move the model",
    "one dict argument keeps the same eval harness (13 holdout)",
], 6.9, 3.1, 5.8, 2.4, size=13)
_, tf = tbox(s, 0.6, 6.1, 12.1, 1.0)
para(tf, 'The idiom axis only exists on graph tasks. On functions, "Python shaped" already is idiomatic, which is why DPO does nothing there and works here.',
     size=13, color=INK, first=True, align=PP_ALIGN.CENTER)

# ----------------------------------------------------------------- 10 GRAPHS & STATS
s = slide(); title(s, "Graphs and stats: behavioral test pass")
bullets(s, [
    "Base eval from **cold start** (quantized Q4/Q8, run normally).",
    "**Base is 0%**: neither model knew Jac before this.",
    "Functions: both reach **~94%**. SFT and DPO tie, outputs stay Python shaped, and DPO does not move them.",
    "Graph: **Qwen pulls ahead** (46/61%) while Gemma stalls (15%).",
], 0.6, 1.4, 5.4, 4.5, size=15)
image(s, "accuracy_compare.png", None, 1.7, w=6.8, center_x=9.5)

# ----------------------------------------------------------------- 11 GRAPHS PT2
s = slide(); title(s, "Graphs and stats, part 2: training curves")
_, tf = tbox(s, 0.6, 1.5, 3.7, 4.5)
para(tf, "Validation loss, training loss, and the holdout learning curve are **nearly identical** for Qwen and Gemma. Same data, same config (LoRA r16, 600 iters, lr 2e-5).",
     size=14, color=INK, first=True, space=12)
para(tf, "The separation is not in loss. It is in idiom on graph tasks, which is the next slide.",
     size=13, color=SOFT)
image(s, "learning_curve_compare.png", 4.5, 1.4, w=4.1)
image(s, "train_loss_compare.png", 8.8, 1.4, w=4.1)
image(s, "val_loss_compare.png", 6.6, 4.3, w=4.1)

# ----------------------------------------------------------------- 12 GRAPH TIER IDIOM
s = slide(); title(s, "Graph tier: where idiom has room to move")
bullets(s, [
    "On graph tasks idiomatic Jac diverges from Python, so the idiom axis is alive.",
    "**SFT** teaches the model to run. **DPO** pushes it the rest of the way to idiomatic.",
    "**Qwen**: SFT 46% → DPO **61%** correct; idiomatic share of correct 83% → **100%**; similarity 0.457 → **0.338** (ref 0.26).",
    "**Gemma**: graph SFT only 15% (~2/13, 0 constructs), so DPO has nothing to move. Idiom acquisition depends on the model.",
], 0.6, 1.4, 6.2, 4.8, size=14)
image(s, "graph_idiom_compare.png", None, 1.9, w=5.6, center_x=10.0)

# ----------------------------------------------------------------- 13 WHAT WE PROVED
s = slide(); title(s, "What we proved")
numbered(s, [
    "**Synthetic, compiler validated data teaches correct Jac.**  0% → **94%** behavioral test pass on a decontaminated holdout.",
    "**Idiom headroom plus DPO teaches idiomatic Jac.**  On graph tasks Qwen3 Coder reaches **100%** idiomatic of correct outputs, up from 83% at SFT.",
    "**The result is a basic Jac coding agent.**  Roughly what Claude Code is, but for Jac instead of Python.",
], 0.9, 1.9, 11.5, 3.4, size=18)
caption(s, "100% synthetic data   •   a free compiler oracle   •   two MoE models that fit in 48 GB", 5.6, size=14, color=SOFT)

# ----------------------------------------------------------------- 14 NEXT UP
s = slide(); title(s, "Next up")
bullets(s, [
    "Mine a single codebase into an **RL loop**: how far is the resulting codebase from what you asked for?",
    "Is it unreasonable to train a model on *one* codebase?",
    "What **eval** fits that setting?",
], 0.6, 1.2, 7.6, 1.9, size=15)
_, tf = tbox(s, 8.5, 1.2, 4.2, 2.0)
para(tf, "**Shipped since:** a live **training dashboard** (Monitor, Train, Ingest) that streams SFT and DPO runs in real time. It is the item that was \"next\" last time.",
     size=12.5, color=INK, first=True)
image(s, "dashboard_wide.png", None, 3.0, w=8.4, center_x=SW/2)
caption(s, "live Monitor: run picker, idiom cards, train and val loss, learning curve, tok/s", 7.0, size=10)

out = os.path.join(HERE, "jaseci_deck.pptx")
prs.save(out)
print("wrote", out, "slides:", len(prs.slides._sldIdLst))

# ---- bounds + content audit (no renderer needed) ----
from pptx.util import Emu as _Emu
chk = Presentation(out)
SWE, SHE = chk.slide_width, chk.slide_height
problems = 0
for i, sl in enumerate(chk.slides, 1):
    ntext = nimg = 0
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            ntext += 1
        if sh.shape_type == 13:  # picture
            nimg += 1
        L, T, W, H = sh.left or 0, sh.top or 0, sh.width or 0, sh.height or 0
        if L < -9525 or T < -9525 or L + W > SWE + 9525 or T + H > SHE + 9525:
            print(f"  ! slide {i} overflow: {sh.shape_type} L={L/914400:.2f} T={T/914400:.2f} R={(L+W)/914400:.2f} B={(T+H)/914400:.2f}")
            problems += 1
    print(f"slide {i:2d}: text_shapes={ntext} images={nimg}")
print("OVERFLOWS:", problems)
